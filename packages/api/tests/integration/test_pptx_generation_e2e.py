"""Spec 16 T12 — pptx_generation end-to-end backend integration test.

Quality bar criterion #2 binary test for the pptx skill (research.md §3.3,
§3.7 pptx rubric). Drives a real :class:`LocalDockerSandbox` through the
full M1a runtime composition path — ``use_skill("pptx_generation")``
activation stages the on-disk ``supplements/`` directory via
:func:`collect_skill_supplements` (D-16-X-7 relative paths), the next
``code_execution`` call writes a 6-slide deck via ``python-pptx``, and the
runtime's ``_persist_produced_file`` callback (api-side
``packages/api/src/persona_api/sandbox/runtime_tool.py::_persist_produced_file``)
copies the produced file into ``<persona_workspace_root>/<owner_id>/
<persona_id>/<filename>`` via :meth:`CodeSandbox.copy_produced_file_to`.

Per the production-fixes phase preconditions (handover.md):

- **D-16-X-6 PATH ordering** (``_BASE_CONTAINER_KWARGS['environment']['PATH']``
  prepends ``/opt/venv/bin``) — verified at module-import time. NO
  ``sys.path.insert`` prelude added; the test snippet imports
  ``from pptx import Presentation`` directly.
- **D-16-X-7 supplements relative-path fix** —
  :func:`collect_skill_supplements` returns ``SandboxFile(path=".skills/...")``
  (relative). Production supplements staging end-to-end works; this test
  uses the real ``pptx_generation`` SkillSpec with its on-disk
  ``supplements/`` directory.

Per Dominant Concern #2 (quality IS the deliverable): the produced file is
parsed via ``pptx.Presentation`` and asserted against every §3.7 row. PARTIAL
visual-surrogate rows are reported as PARTIAL in the state.md scorecard, not
laundered to PASS.

Per T01 audit A1: ``_ScriptedBackend`` is module-private per the api per-file
convention (three precedents). The pattern mirrors
``packages/runtime/tests/_fakes.py:50``.

Markers: ``integration`` AND ``docker``. Skips cleanly when Docker is
unreachable (per A1 convention).
"""

# ruff: noqa: ANN401, ARG002, SLF001 — test doubles with intentionally loose
# sigs (ANN401/ARG002); sandbox host_out introspection for cleanup (SLF001).

from __future__ import annotations

import shutil
from pathlib import Path
from typing import TYPE_CHECKING, Any

import pytest
from persona.backends.types import ChatResponse, TokenUsage
from persona.sandbox.local_docker import (
    DEFAULT_IMAGE,
    LocalDockerSandbox,
    is_docker_available,
)
from persona.sandbox.result import NetworkPolicy, ResourceLimits, SandboxFile
from persona.sandbox.tool import make_code_execution_tool
from persona.schema.persona import Persona, PersonaIdentity
from persona.schema.skills import SkillSpec
from persona.schema.tools import ToolCall
from persona.skills import (
    SkillInjector,
    SkillScanner,
    make_use_skill_tool,
)
from persona.tools import Toolbox
from persona_runtime.agentic.loop import AgenticLoop
from persona_runtime.agentic.run import RunStatus
from persona_runtime.agentic.step import StepType
from persona_runtime.prompt import PromptBuilder
from persona_runtime.router import Router
from persona_runtime.tier import TierConfig, TierRegistry

if TYPE_CHECKING:
    from collections.abc import AsyncIterator

    from persona.backends.types import StreamChunk, ToolSpec
    from persona.schema.conversation import ConversationMessage


pytestmark = [pytest.mark.integration, pytest.mark.docker]


# Built-in skills root — point the scanner at the on-disk packs so the real
# ``pptx_generation`` SKILL.md + supplements directory are exercised
# end-to-end (M1a + D-16-X-7 verification path).
_BUILTIN_ROOT = (
    Path(__file__).parent.parent.parent.parent / "core" / "src" / "persona" / "skills" / "builtin"
).resolve()

# Inspection artifact directory (D-16-X-3 — gitignored). The test writes a
# copy of the produced file here so the operator can re-open it and the
# evidence is reproducible by re-running the test. ``parents[4]`` walks up
# ``integration → tests → api → packages → <repo-root>`` from this file's
# location; ``parents[3]`` would land at ``packages/`` (the wrong root).
_INSPECTION_DIR = (
    Path(__file__).resolve().parents[4] / "docs" / "specs" / "phase2" / "spec_16" / "inspection"
).resolve()

# §3.3 representative task — encoded as the in-sandbox python-pptx code the
# scripted backend issues on the code_execution call. The snippet models
# what an agent that has read the on-disk SKILL.md + supplements would write
# idiomatically: one ``Presentation()`` master, layout-per-intent
# selection, speaker notes on slides 3/4/6 via
# ``slide.notes_slide.notes_text_frame.text``, a matplotlib-rendered bar
# chart embedded on slide 4 via ``shapes.add_picture``, descriptive
# filename, valid output path under ``/workspace/out/``.
_PPTX_CODE = """\
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE = prs.slide_layouts[0]
TITLE_CONTENT = prs.slide_layouts[1]

# Slide 1 — title.
s = prs.slides.add_slide(TITLE)
s.shapes.title.text = "Tenant-protection workshop"
sub = s.placeholders[1]
sub.text = "Quarterly briefing - June 2026"

# Slide 2 — agenda.
s = prs.slides.add_slide(TITLE_CONTENT)
s.shapes.title.text = "Agenda"
body = s.placeholders[1].text_frame
body.text = "Pre-2024 rules"
for line in (
    "Post-2024 rules and notice periods",
    "Embedded comparison chart",
    "Q and A",
    "Summary",
):
    p = body.add_paragraph()
    p.text = line

# Slide 3 — pre-2024 rules + speaker notes (hard feature).
s = prs.slides.add_slide(TITLE_CONTENT)
s.shapes.title.text = "Pre-2024 rules"
body = s.placeholders[1].text_frame
body.text = "Three months notice for ordinary increases"
for line in (
    "Tenant may dispute within 30 days",
    "Index-linked increases capped at CPI",
):
    p = body.add_paragraph()
    p.text = line
s.notes_slide.notes_text_frame.text = (
    "Walk through the three pre-2024 protection clauses. "
    "Emphasise that the three-month notice was insufficient for "
    "vulnerable tenants in the 2022 audit. The comparison chart "
    "is on the next slide."
)

# Slide 4 — post-2024 rules + embedded bar chart of notice-period months
# by tenure-year + speaker notes (hard feature).
chart_path = Path("/workspace/out/notice-periods.png")
fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
tenure_years = ["<1y", "1-3y", "3-5y", "5y+"]
notice_months = [3, 4, 5, 6]
ax.bar(tenure_years, notice_months, color="#1f4e79")
ax.set_title("Notice-period months by tenure-year (post-2024)")
ax.set_xlabel("Tenure")
ax.set_ylabel("Months notice")
for i, v in enumerate(notice_months):
    ax.text(i, v + 0.1, str(v), ha="center")
fig.tight_layout()
fig.savefig(chart_path, format="png")
plt.close(fig)

s = prs.slides.add_slide(TITLE_CONTENT)
s.shapes.title.text = "Post-2024 rules"
body = s.placeholders[1].text_frame
body.text = "Six-month notice for tenants of five years and longer"
body.add_paragraph().text = "Tiered notice by tenure shown right"
s.shapes.add_picture(
    str(chart_path),
    left=Inches(6.5),
    top=Inches(1.8),
    width=Inches(6.0),
    height=Inches(3.4),
)
s.notes_slide.notes_text_frame.text = (
    "The bar chart shows the tiered notice schedule by tenure-year. "
    "Note the doubling from three to six months for tenants of five "
    "years and longer. Pause for the room to study the chart."
)

# Slide 5 — Q and A prompt.
s = prs.slides.add_slide(TITLE_CONTENT)
s.shapes.title.text = "Questions and discussion"
body = s.placeholders[1].text_frame
body.text = "What scenarios does the new schedule not cover?"
body.add_paragraph().text = "How will municipal tribunals handle disputes?"
body.add_paragraph().text = "Open floor"

# Slide 6 — summary + speaker notes (hard feature).
s = prs.slides.add_slide(TITLE_CONTENT)
s.shapes.title.text = "Summary"
body = s.placeholders[1].text_frame
body.text = "Pre-2024: uniform three-month notice"
for line in (
    "Post-2024: tiered notice by tenure (three to six months)",
    "Tenants of five years and longer gain the most protection",
    "Tribunals retain dispute jurisdiction",
):
    p = body.add_paragraph()
    p.text = line
s.notes_slide.notes_text_frame.text = (
    "Wrap up by restating the three takeaways. Invite participants to "
    "circulate the comparison chart to their members. Thank the room."
)

out_path = Path("/workspace/out/tenant-protection-workshop.pptx")
prs.save(out_path)
print(
    f"wrote {out_path.name} "
    f"({out_path.stat().st_size} bytes, {len(prs.slides)} slides)"
)
"""


# ----- module-private scripted backend (T01 audit A1 per-file convention) ---


class _ScriptedBackend:
    """Replays a fixed sequence of :class:`ChatResponse` objects.

    Mirrors ``packages/runtime/tests/_fakes.py:50`` per T01 audit A1.
    Only the ``chat()`` non-streaming surface is implemented because the
    :class:`AgenticLoop` drives via ``chat()``; ``chat_stream`` raises if
    the loop ever calls it (defensive — the AgenticLoop never does).
    """

    def __init__(
        self,
        chat_script: list[ChatResponse],
        *,
        provider_name: str = "anthropic",
        model_name: str = "claude-sonnet-4-6",
    ) -> None:
        self._chat_script = list(chat_script)
        self._chat_index = 0
        self._provider_name = provider_name
        self._model_name = model_name
        self.chat_calls = 0

    @property
    def provider_name(self) -> str:
        return self._provider_name

    @property
    def model_name(self) -> str:
        return self._model_name

    @property
    def supports_native_tools(self) -> bool:
        return False

    @property
    def supports_vision(self) -> bool:
        return False

    async def chat(
        self,
        messages: list[ConversationMessage],
        *,
        tools: list[ToolSpec] | None = None,
        temperature: float = 0.0,
        max_tokens: int = 4096,
        stop: list[str] | None = None,
    ) -> ChatResponse:
        self.chat_calls += 1
        if self._chat_index < len(self._chat_script):
            response = self._chat_script[self._chat_index]
            self._chat_index += 1
            return response
        # Defensive empty final — the loop should have terminated already.
        return ChatResponse(
            content="[FINAL] ",
            tool_calls=[],
            usage=TokenUsage(prompt_tokens=1, completion_tokens=0, total_tokens=1),
            model=self._model_name,
            provider=self._provider_name,
            latency_ms=0.0,
        )

    async def chat_stream(
        self,
        messages: list[ConversationMessage],
        *,
        tools: list[ToolSpec] | None = None,
        temperature: float = 0.0,
        max_tokens: int = 4096,
        stop: list[str] | None = None,
    ) -> AsyncIterator[StreamChunk]:
        raise NotImplementedError("AgenticLoop drives chat(), not chat_stream()")
        yield  # pragma: no cover — unreachable, satisfies async-generator typing


# ----- helpers --------------------------------------------------------------


def _image_present(tag: str) -> bool:
    """True if ``tag`` is locally available on the Docker daemon."""
    try:
        import docker
        from docker.errors import DockerException, ImageNotFound
    except ImportError:
        return False
    try:
        client = docker.from_env()
        try:
            client.images.get(tag)
        except ImageNotFound:
            return False
        finally:
            client.close()
    except DockerException:
        return False
    return True


def _docker_skip_reason() -> str | None:
    """Return a skip reason if Docker isn't ready, else ``None``."""
    if not is_docker_available():
        return (
            "Docker daemon unreachable; T12 needs LocalDockerSandbox to drive "
            "the real pptx_generation production path. Start Docker and rerun."
        )
    if not _image_present(DEFAULT_IMAGE):
        return (
            f"Sandbox image {DEFAULT_IMAGE!r} not built. Build via the Spec 12 "
            "T06 Dockerfile (or pull) before running T12."
        )
    return None


def _persona() -> Persona:
    """Build the minimal persona that activates ``pptx_generation``."""
    return Persona(
        persona_id="astrid_t12_pptx",
        identity=PersonaIdentity(
            name="Astrid",
            role="Norwegian tenancy workshop facilitator",
            background=("Produces presentation materials for tenant-protection workshops."),
            language_default="nb",
            constraints=["Do not give binding legal advice."],
        ),
        skills=["pptx_generation"],
        tools=["use_skill", "code_execution"],
    )


def _empty_stores() -> dict[str, Any]:
    """Minimal in-memory MemoryStore doubles — T12 asserts on files, not memory."""

    class _MemStore:
        def write(self, *_a: Any, **_kw: Any) -> None: ...

        def query(self, *_a: Any, **_kw: Any) -> list[Any]:
            return []

        def get_all(self, *_a: Any, **_kw: Any) -> list[Any]:
            return []

        def delete(self, *_a: Any, **_kw: Any) -> None: ...

        def remove_documents(self, *_a: Any, **_kw: Any) -> None: ...

        def history(self, *_a: Any, **_kw: Any) -> list[Any]:
            return []

        def rollback(self, *_a: Any, **_kw: Any) -> None: ...

    return {
        "identity": _MemStore(),
        "self_facts": _MemStore(),
        "worldview": _MemStore(),
        "episodic": _MemStore(),
    }


def _scan(skill_names: list[str]) -> list[SkillSpec]:
    """Scan the named built-in skills and assert all resolved."""
    scanner = SkillScanner([_BUILTIN_ROOT])
    specs = scanner.scan(skill_names)
    assert len(specs) == len(skill_names), (
        f"scanner returned {len(specs)} specs for {skill_names!r}; "
        "expected all to resolve from the built-in path"
    )
    return specs


def _resp(content: str = "", *, tool_calls: list[ToolCall] | None = None) -> ChatResponse:
    return ChatResponse(
        content=content,
        tool_calls=tool_calls or [],
        usage=TokenUsage(prompt_tokens=10, completion_tokens=5, total_tokens=15),
        model="claude-sonnet-4-6",
        provider="anthropic",
        latency_ms=1.0,
    )


def _build_loop(
    *,
    persona: Persona,
    backend: _ScriptedBackend,
    toolbox: Toolbox,
    scanned_skills: list[SkillSpec],
    max_steps: int = 20,
) -> AgenticLoop:
    """Wire an :class:`AgenticLoop` over the supplied collaborators."""
    from persona.backends import BackendConfig

    dummy_cfg = BackendConfig(provider="anthropic", model="m", api_key=None)
    registry = TierRegistry(
        {
            "frontier": TierConfig(name="frontier", backend_config=dummy_cfg),
            "mid": TierConfig(name="mid", backend_config=dummy_cfg),
            "small": TierConfig(name="small", backend_config=dummy_cfg),
        }
    )
    registry._cache = {  # type: ignore[assignment]
        "frontier": backend,  # type: ignore[dict-item]
        "mid": backend,  # type: ignore[dict-item]
        "small": backend,  # type: ignore[dict-item]
    }
    return AgenticLoop(
        persona=persona,
        stores=_empty_stores(),  # type: ignore[arg-type]
        toolbox=toolbox,
        skill_injector=SkillInjector(),
        scanned_skills=scanned_skills,
        prompt_builder=PromptBuilder(),
        router=Router(),
        tier_registry=registry,
        max_steps=max_steps,
    )


# ----- the production-verification test -------------------------------------


class TestT12PptxEndToEnd:
    """T12 — production-verification path for the ``pptx_generation`` skill.

    Drives the full M1a runtime composition + bytes-persistence pipeline
    (D-16-X-6 PATH ordering, D-16-X-7 supplements-relative-path,
    D-12-X-read-produced-file copy callback, D-17-X-bytes-persistence
    workspace persistence) and parses the resulting deck via
    ``pptx.Presentation`` to score the §3.7 pptx rubric.
    """

    @pytest.mark.asyncio
    async def test_pptx_generation_produces_workshop_deck_end_to_end(
        self,
        tmp_path: Path,
    ) -> None:
        """The persona activates ``pptx_generation``; the code_execution call
        writes a 6-slide deck satisfying every §3.7 row.
        """
        skip_reason = _docker_skip_reason()
        if skip_reason:
            pytest.skip(skip_reason)

        # 1) Pre-condition guards (production-fixes phase preconditions per
        # handover.md). Read post-fix file content directly — if these
        # invariants regress the test fails fast with a clear signal rather
        # than executing the snippet and looking like a Docker / library
        # bug.

        # D-16-X-6 PATH ordering: ``/opt/venv/bin`` must be FIRST in the
        # _BASE_CONTAINER_KWARGS["environment"]["PATH"] override so the
        # image's venv-installed python-pptx + matplotlib are reachable
        # natively (no sys.path.insert workaround in the snippet).
        from persona.sandbox.local_docker import _BASE_CONTAINER_KWARGS

        path_value = _BASE_CONTAINER_KWARGS["environment"]["PATH"]
        path_segments = path_value.split(":")
        assert path_segments[0] == "/opt/venv/bin", (
            f"D-16-X-6 regressed: PATH={path_value!r} (expected /opt/venv/bin first); "
            "sys.path.insert prelude would be needed if this is not first."
        )

        # D-16-X-7 supplements-relative-path: collect_skill_supplements must
        # emit relative ``.skills/<name>/...`` paths (NEVER absolute) so the
        # consumer-side join lands at <host_in>/.skills/... rather than
        # writing to the host's /workspace/in/ (which fails on macOS with
        # OSError [Errno 30] Read-only file system).
        from persona.skills import collect_skill_supplements

        scanned = _scan(["pptx_generation"])
        pptx_spec = next(s for s in scanned if s.name == "pptx_generation")
        supplement_files = collect_skill_supplements(pptx_spec)
        assert len(supplement_files) >= 1, (
            "pptx_generation must ship on-disk supplements/ markdown files "
            "(T04 close-gate documented 3 supplements)"
        )
        for sf in supplement_files:
            assert not sf.path.startswith("/"), (
                f"D-16-X-7 regressed: collect_skill_supplements emitted absolute "
                f"path {sf.path!r}; should be .skills/... relative."
            )
            assert sf.path.startswith(".skills/pptx_generation/supplements/"), (
                f"unexpected supplement path shape: {sf.path!r}"
            )

        # 2) Build a real LocalDockerSandbox + composition root that mirrors
        # the api production wiring: deferred_input_files_provider (the
        # M1a drain) + produced_file_persister (the D-17-X-bytes-persistence
        # copy callback). The persona workspace root lives under tmp_path
        # so the test is hermetic; the persister copies into
        # <persona_workspace_root>/<owner_id>/<persona_id>/<filename>.
        owner_id = "owner_t12"
        persona = _persona()
        workspace_root = tmp_path / "wsp"
        persona_workspace = workspace_root / owner_id / (persona.persona_id or "astrid")
        persona_workspace.mkdir(parents=True, exist_ok=True)

        sandbox = LocalDockerSandbox(workspace_root=tmp_path / "sbx")
        session_id = f"{owner_id}:conv_t12_pptx"

        # Shared list[SandboxFile] holder — same identity as the loop's
        # public ``deferred_input_files`` attribute and the tool's
        # drain-and-clear provider closure. Mirrors
        # ``persona_api.services.runtime_factory.RuntimeFactory``'s
        # composition root (D-16-2-state-location).
        deferred_holder: list[SandboxFile] = []

        def _drain_and_clear() -> list[SandboxFile]:
            snapshot = list(deferred_holder)
            deferred_holder.clear()
            return snapshot

        async def _persist_produced_file(sid: str, ref: str) -> None:
            """Copy a produced file from the sandbox session to the persona
            workspace via the shipped :meth:`copy_produced_file_to` method.

            Mirrors the production callback at
            ``packages/api/src/persona_api/sandbox/runtime_tool.py:216``.
            """
            await sandbox.copy_produced_file_to(sid, ref, persona_workspace / ref)

        produced_path: Path | None = None
        try:
            await sandbox.create_session(
                session_id,
                limits=ResourceLimits(),
                network=NetworkPolicy(),
            )
            code_exec = make_code_execution_tool(
                sandbox,
                persona_id=persona.persona_id,
                session_id_provider=lambda: session_id,
                deferred_input_files_provider=_drain_and_clear,
                produced_file_persister=_persist_produced_file,
            )
            use_skill = make_use_skill_tool(scanned)
            toolbox = Toolbox(
                [use_skill, code_exec],
                allow_list=["use_skill", "code_execution"],
            )

            # 3) Scripted backend — (a) activate pptx_generation; (b) code_execution
            # writes the 6-slide deck; (c) FINAL.
            chat_script = [
                _resp(
                    tool_calls=[
                        ToolCall(
                            name="use_skill",
                            args={"skill_name": "pptx_generation"},
                            call_id="us1",
                        )
                    ]
                ),
                _resp(
                    tool_calls=[
                        ToolCall(
                            name="code_execution",
                            args={"code": _PPTX_CODE},
                            call_id="ce1",
                        )
                    ]
                ),
                _resp("[FINAL] Wrote /workspace/out/tenant-protection-workshop.pptx"),
            ]
            backend = _ScriptedBackend(chat_script)
            loop = _build_loop(
                persona=persona,
                backend=backend,
                toolbox=toolbox,
                scanned_skills=scanned,
            )
            # Bind the loop's public attribute to the shared holder so the
            # use_skill intercept's extend(...) and the tool's drain-and-clear
            # see the same list identity (D-16-2-state-location).
            loop.deferred_input_files = deferred_holder

            run = await loop.run("Produce the 6-slide tenant-protection workshop deck")

            # Run terminated cleanly.
            assert run.status is RunStatus.COMPLETED, (
                f"run did not complete: status={run.status} error={run.error!r}"
            )

            # Step shape: use_skill, code_execution, FINAL.
            step_types = [s.type for s in run.steps]
            assert step_types == [
                StepType.TOOL_CALL,
                StepType.TOOL_CALL,
                StepType.FINAL,
            ], f"unexpected step type sequence: {step_types}"

            # 4) Sandbox-boundary assertion (T01 audit A2 / D-16-X-5
            # same-session reachability) — the tool result mirrors
            # ExecutionResult.produced_files into ToolResult.data
            # ['produced_files']. Required for #1 (file produced) + #5
            # (sandbox-side library import succeeded).
            code_step = next(
                s
                for s in run.steps
                if s.type is StepType.TOOL_CALL
                and any(c.name == "code_execution" for c in s.tool_calls)
            )
            code_result = code_step.results[0]
            assert code_result.is_error is False, (
                f"code_execution must succeed; got error: {code_result.content!r}"
            )
            assert code_result.data is not None
            produced = code_result.data.get("produced_files") or []
            pptx_entries = [p for p in produced if str(p.get("path", "")).endswith(".pptx")]
            assert len(pptx_entries) == 1, (
                f"expected exactly one .pptx in produced_files; got {produced}"
            )
            assert pptx_entries[0]["size_bytes"] > 0

            # 5) Workspace persistence assertion (D-17-X-bytes-persistence) —
            # the file landed at <persona_workspace>/<filename> via the
            # produced_file_persister callback (criterion #3 — produced-files
            # contract → workspace, #11 — RLS via per-tenant workspace).
            expected_path = persona_workspace / "tenant-protection-workshop.pptx"
            assert expected_path.is_file(), (
                f"produced .pptx not persisted to persona workspace at "
                f"{expected_path}; produced={produced}"
            )
            assert expected_path.stat().st_size > 5_000, (
                f".pptx is suspiciously small ({expected_path.stat().st_size} "
                "bytes); a valid minimal deck is ~30 KB."
            )
            produced_path = expected_path

            # 6) Parse the produced deck and score the §3.7 rubric.
            from pptx import Presentation

            prs = Presentation(str(produced_path))
            # ``prs.slides`` does not support slicing in python-pptx 1.0.2
            # (``__getitem__`` only accepts int per a slide-id lookup);
            # materialise the slide list once and reuse for the §3.7 rubric.
            slides_list = list(prs.slides)

            # The §3.7 programmatic asserts — exact shapes from research §3.7
            # pptx rubric. Track PASS / PARTIAL outcomes for the state.md
            # scorecard at the end of the test.
            asserts_passed = 0
            asserts_partial = 0
            visual_surrogates = 0

            # Row 1 — slide-master shared. Note: ``SlideMaster`` instances are
            # not hashable in python-pptx 1.0.2, so the §3.7 surface
            # ``len({s.slide_layout.slide_master for s in prs.slides}) == 1``
            # cannot be used verbatim; substitute identity-based
            # de-duplication via ``id()``.
            master_ids = {id(s.slide_layout.slide_master) for s in slides_list}
            assert len(master_ids) == 1, (
                f"§3.7 row 1 (slide-master shared) failed: {len(master_ids)} "
                "distinct master object identities"
            )
            asserts_passed += 1

            # Row 2 — layout chosen explicitly (slide 0 is title, the rest are
            # non-Blank). The §3.7 surface allows slide 0 to be title; check
            # slides[1:] for non-Blank.
            for idx, slide in enumerate(slides_list[1:], start=1):
                assert slide.slide_layout.name != "Blank", (
                    f"§3.7 row 2 failed: slide {idx} uses Blank layout"
                )
            asserts_passed += 1

            # Row 3 — speaker notes on slides 3/4/6 (indices 2/3/5).
            for idx in (2, 3, 5):
                notes_text = slides_list[idx].notes_slide.notes_text_frame.text.strip()
                assert notes_text, f"§3.7 row 3 failed: slide index {idx} has empty speaker notes"
            asserts_passed += 1

            # Row 4 — title ≥28pt, body ≥18pt. The shipped python-pptx default
            # master honours this without an explicit run-level override, so
            # this is a structural / visual row. Programmatic surrogate:
            # walk every text frame and accumulate explicit run font-size
            # values; assert no explicit override REDUCES below the minimum.
            # Implicit-default text frames remain visual-only.
            explicit_sizes: list[tuple[str, int]] = []
            for slide in slides_list:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    placeholder_kind = "body"
                    if shape == slide.shapes.title:
                        placeholder_kind = "title"
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                explicit_sizes.append((placeholder_kind, run.font.size.pt))
            # All explicit overrides — if any — must respect the §3.7 floor.
            for kind, size in explicit_sizes:
                if kind == "title":
                    assert size >= 28, (
                        f"§3.7 row 4 (title ≥28pt) failed via explicit run override: {size}pt"
                    )
                else:
                    assert size >= 18, (
                        f"§3.7 row 4 (body ≥18pt) failed via explicit run override: {size}pt"
                    )
            # The structural side passes (no run drops below floor); the
            # full visual confirmation requires opening in PowerPoint /
            # LibreOffice Impress, so this is PARTIAL by §3.7's "Visual"
            # note. Reported honestly in the scorecard.
            asserts_partial += 1
            visual_surrogates += 1

            # Row 5 — no off-slide overflow. Purely visual per §3.7.
            # Programmatic surrogate: per-slide shape left+width and
            # top+height fit inside prs.slide_width / prs.slide_height.
            slide_w_emu = prs.slide_width
            slide_h_emu = prs.slide_height
            for slide_idx, slide in enumerate(slides_list):
                for shape in slide.shapes:
                    # Skip placeholders without explicit coordinates.
                    if shape.left is None or shape.top is None:
                        continue
                    if shape.width is None or shape.height is None:
                        continue
                    right = shape.left + shape.width
                    bottom = shape.top + shape.height
                    # Allow 5% slack — the master may legitimately overhang
                    # by a fraction due to padding.
                    assert right <= int(slide_w_emu * 1.05), (
                        f"§3.7 row 5 surrogate failed: slide {slide_idx} shape overflows right edge"
                    )
                    assert bottom <= int(slide_h_emu * 1.05), (
                        f"§3.7 row 5 surrogate failed: slide {slide_idx} "
                        f"shape overflows bottom edge"
                    )
            # Best-effort programmatic surrogate; full visual verification
            # requires opening in PowerPoint.
            asserts_partial += 1
            visual_surrogates += 1

            # Row 6 — no mixed fonts per slide (≤ 2 distinct explicit
            # font-name values per slide). Without explicit per-run font
            # overrides this is implicitly satisfied via the master.
            for slide_idx, slide in enumerate(slides_list):
                font_names = {
                    r.font.name
                    for sh in slide.shapes
                    if sh.has_text_frame
                    for p in sh.text_frame.paragraphs
                    for r in p.runs
                    if r.font.name
                }
                assert len(font_names) <= 2, (
                    f"§3.7 row 6 failed: slide {slide_idx} has "
                    f"{len(font_names)} distinct fonts {font_names}"
                )
            asserts_passed += 1

            # Row 7 — chart embed renders on slide 4 (index 3). §3.7's
            # canonical assert: ``any(sh.shape_type == 13 for sh in
            # prs.slides[3].shapes)`` (13 = PICTURE). The task brief widens
            # to (13, 3) to also accept native CHART (3) since the
            # representative task is ambiguous on raster vs native.
            slide_4_shape_types = {sh.shape_type for sh in slides_list[3].shapes}
            assert any(st in (13, 3) for st in slide_4_shape_types), (
                f"§3.7 row 7 failed: slide 4 shape types {slide_4_shape_types} "
                "include neither PICTURE (13) nor CHART (3)"
            )
            asserts_passed += 1

            # Row 8 — slide count = 6 (the representative task contract).
            assert len(slides_list) == 6, (
                f"§3.7 row 8 failed: got {len(slides_list)} slides; expected 6"
            )
            asserts_passed += 1

            # 7) Save the inspection artifact (D-16-X-3 — gitignored; the
            # operator can re-open or the scorecard references it by
            # filename). ``mkdir(parents=True, exist_ok=True)`` first.
            _INSPECTION_DIR.mkdir(parents=True, exist_ok=True)
            inspection_path = _INSPECTION_DIR / "pptx_sample.pptx"
            shutil.copy(produced_path, inspection_path)
            assert inspection_path.is_file()

            # 8) Final sanity — record the counts on the test class so the
            # scorecard reader sees the numbers. The state.md scorecard is
            # appended in a separate write step at the close of T12.
            assert asserts_passed >= 5, (
                f"expected at least 5 fully-programmatic PASS asserts; got {asserts_passed}"
            )
            assert asserts_partial >= 1, (
                f"expected at least 1 PARTIAL surrogate; got {asserts_partial}"
            )
            assert visual_surrogates >= 1, (
                f"expected at least 1 visual surrogate; got {visual_surrogates}"
            )
        finally:
            await sandbox.aclose()
